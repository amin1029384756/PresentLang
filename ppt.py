import pptx
from pptx.util import Inches
from pptx.dml.color import RGBColor

# Define a mapping dictionary for slide layouts
SLIDE_LAYOUT_MAP = {
    'intro': 5,  # Title Only layout (for centered title and subtitle)
    'default': 2,  # Title and Content layout
    'blank': 6,
    'title': 1,
    'section_header': 3,
    'title_and_two_content': 4,
    # Add more layout mappings as needed
}

# Define a default color for slides
DEFAULT_SLIDE_COLOR = RGBColor(0, 0, 0)  # Black

def create_slides_presentation(presentation_details, scenes):
    # Create a new PowerPoint presentation
    presentation = pptx.Presentation()

    # Set the presentation metadata
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    # Iterate over each scene and create a slide
    for scene_name, scene_details in scenes.items():
        # Get the slide layout index from the mapping dictionary
        layout_index = SLIDE_LAYOUT_MAP.get(scene_details['type'], 0)

        # Create a new slide with the specified layout
        slide_layout = presentation.slide_layouts[layout_index]
        slide = presentation.slides.add_slide(slide_layout)
        print(scene_details['subtitle'])
        print(scene_details)
        # Handle different scene types
        if scene_details['type'] == 'intro':
            title_shape = slide.shapes.title
            title_shape.text = scene_details['title']
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(1))
            subtitle_box.text_frame.text = scene_details['subtitle']

        elif scene_details['type'] == 'default':
            # Set the title in the top left
            slide.shapes.title.text = scene_details['title']

            # Add text content to the content placeholder
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = scene_details['text']

            imgpath = scene_details['imagepath']
            if imgpath != "":
                slide.shapes.add_picture(imgpath, Inches(8), Inches(2), width=Inches(8), height=Inches(4.5))
            else:
                continue
        # Set the slide background and theme
        slide.background.fill.solid()
        try:
            slide.background.fill.fore_color.rgb = RGBColor.from_string(scene_details['color'])
        except ValueError:
            slide.background.fill.fore_color.rgb = DEFAULT_SLIDE_COLOR

    # Save the presentation
    presentation.save('output.pptx')